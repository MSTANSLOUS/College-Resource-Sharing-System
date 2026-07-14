import os
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import io
from PIL import Image, ImageDraw, ImageFont  # PyMuPDF
from flask import Blueprint, render_template, request, redirect, current_app, flash, url_for, abort, jsonify
from flask_login import login_required, current_user, login_user, logout_user
from werkzeug.utils import secure_filename
from werkzeug.security import check_password_hash, generate_password_hash
from app.models import db, Program, Resource, User, Log, Module, TransferRequest
from datetime import date, datetime
from app import send_email, socketio
from sqlalchemy import func

core = Blueprint('core', __name__)

def create_log(action, details=None):
    log_entry = Log(
        user_id=current_user.id if current_user.is_authenticated else None,
        action=action,
        details=details
    )
    db.session.add(log_entry)
    db.session.commit()


# Add this helper function at the top
def paginate_query(query, page=1, per_page=10):
    total = query.count()
    items = query.offset((page - 1) * per_page).limit(per_page).all()
    return {
        'items': items,
        'total': total,
        'page': page,
        'per_page': per_page,
        'pages': (total + per_page - 1) // per_page if total > 0 else 1
    }


# ─── Paginated Table Fragments ───

@core.route('/admin/table/reps')
@login_required
def admin_reps_table():
    if not current_user.is_admin:
        abort(403)
    page = request.args.get('page', 1, type=int)
    per_page = 10
    users = User.query.filter_by(is_class_rep=True, is_admin=False)
    paginated = paginate_query(users, page, per_page)
    return render_template('core/admin/_reps_table_body.html', 
                           users=paginated['items'],
                           page=paginated['page'],
                           pages=paginated['pages'])



@core.route('/admin/table/students')
@login_required
def admin_students_table():
    if not current_user.is_admin:
        abort(403)
    page = request.args.get('page', 1, type=int)
    per_page = 10
    year_filter = request.args.get('year', type=int)
    sem_filter = request.args.get('semester', type=int)
    query = User.query.filter(User.is_class_rep == False, User.is_admin == False)
    if year_filter:
        query = query.filter_by(year=year_filter)
    if sem_filter:
        query = query.filter_by(semester=sem_filter)
    paginated = paginate_query(query, page, per_page)
    return render_template('core/admin/_students_table_body.html',
                           users=paginated['items'],
                           page=paginated['page'],
                           pages=paginated['pages'],
                           current_year=year_filter,
                           current_sem=sem_filter)







@core.route('/admin/table/modules')
@login_required
def admin_modules_table():
    if not current_user.is_admin:
        abort(403)
    page = request.args.get('page', 1, type=int)
    per_page = 10
    query = Module.query
    paginated = paginate_query(query, page, per_page)
    return render_template('core/admin/_modules_table_body.html',
                           modules=paginated['items'],
                           page=paginated['page'],
                           pages=paginated['pages'])

@core.route('/admin/table/users')
@login_required
def admin_users_table():
    if not current_user.is_admin:
        abort(403)
    page = request.args.get('page', 1, type=int)
    per_page = 10
    query = User.query
    paginated = paginate_query(query, page, per_page)
    return render_template('core/admin/_users_table_body.html',
                           users=paginated['items'],
                           page=paginated['page'],
                           pages=paginated['pages'],
                           now=datetime.utcnow())

@core.route('/admin/table/logs')
@login_required
def admin_logs_table():
    if not current_user.is_admin:
        abort(403)
    page = request.args.get('page', 1, type=int)
    per_page = 10
    action_filter = request.args.get('action')
    user_search = request.args.get('user_search', '').strip()
    log_query = Log.query
    if action_filter:
        log_query = log_query.filter(Log.action == action_filter)
    if user_search:
        users_matching = User.query.filter(
            (User.full_name.ilike(f'%{user_search}%')) | (User.email.ilike(f'%{user_search}%'))
        ).all()
        user_ids = [u.id for u in users_matching]
        if user_ids:
            log_query = log_query.filter(Log.user_id.in_(user_ids))
        else:
            log_query = log_query.filter(False)
    paginated = paginate_query(log_query.order_by(Log.created_at.desc()), page, per_page)
    return render_template('core/admin/_logs_table_body.html',
                           logs=paginated['items'],
                           page=paginated['page'],
                           pages=paginated['pages'])

@core.route('/dashboard/recent-activity')
@login_required
def recent_activity():
    page = request.args.get('page', 1, type=int)
    per_page = 10
    # For now, we'll show recent logs for the user? Actually we need a recent_activity model.
    # We'll use Logs filtered by user_id
    query = Log.query.filter_by(user_id=current_user.id).order_by(Log.created_at.desc())
    paginated = paginate_query(query, page, per_page)
    activities = []
    for log in paginated['items']:
        activities.append({
            'description': f"{log.action} - {log.details or ''}",
            'time_ago': log.created_at.strftime('%H:%M %d %b %Y'),
            'icon': 'fa-circle-info'
        })
    return render_template('core/student/_recent_activity_list.html',
                           activities=activities,
                           page=paginated['page'],
                           pages=paginated['pages'])







def notify_class_rep(student, user_actions=None):
    rep = User.query.filter_by(
        is_class_rep=True,
        campus=student.campus,
        program_id=student.program_id,
        year=student.year,
        semester=student.semester
    ).first()
    if rep and rep.email:
        subjects = f"Account Creation"
        message_body = f"Hi class admin, {student.full_name} created an account and is waiting for your approval."
        send_email(recipients=rep.email, subject=subjects, massage_body=message_body)







# ─── DASHBOARD ───
@core.route('/dashboard')
@login_required
def dashboard():
    pending_students = []
    if current_user.is_class_rep:
        pending_students = User.query.filter_by(
            is_approved=False,
            campus=current_user.campus,
            program_id=current_user.program_id,
            year=current_user.year,
            semester=current_user.semester
        ).all()

    recent_uploads = Resource.query.filter_by(
        campus=current_user.campus,
        target_year=current_user.year,
        target_semester=current_user.semester
    ).join(Resource.programs).filter(Program.id == current_user.program_id).order_by(Resource.id.desc()).limit(3).all()

    return render_template('core/student/dashboard.html',
                           pending_students=pending_students,
                           recent_uploads=recent_uploads)







@core.route('/live/dashboard-data')
@login_required
def live_dashboard_data():
    recent_uploads = Resource.query.filter_by(
        campus=current_user.campus,
        target_year=current_user.year,
        target_semester=current_user.semester
    ).join(Resource.programs).filter(Program.id == current_user.program_id).order_by(Resource.id.desc()).limit(3).all()

    pending_students = []
    if current_user.is_class_rep:
        pending_students = User.query.filter_by(
            is_approved=False,
            campus=current_user.campus,
            program_id=current_user.program_id,
            year=current_user.year,
            semester=current_user.semester
        ).all()

    uploads_data = []
    for r in recent_uploads:
        uploads_data.append({
            'id': r.id,
            'title': r.title,
            'module_name': r.module.name if r.module else '',
            'academic_year': r.academic_year,
            'file_path': r.file_path,
            'uploader_name': r.uploader.full_name if r.uploader else '',
        })

    approvals_data = []
    for s in pending_students:
        approvals_data.append({
            'id': s.id,
            'full_name': s.full_name,
            'email': s.email,
            'campus': s.campus,
            'year': s.year,
            'semester': s.semester,
        })

    return jsonify({
        'recent_uploads': uploads_data,
        'pending_approvals': approvals_data,
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    })







# ─── PROFILE ───
@core.route('/profile', methods=['GET', 'POST'])
@login_required
def profile():
    if request.method == 'POST':
        old_name = current_user.full_name
        current_user.full_name = request.form.get('fullname')
        current_user.student_id = request.form.get('student_id')
        current_user.phone_number = request.form.get('phone')

        new_year = request.form.get('year')
        new_sem = request.form.get('semester')
        if new_year: current_user.year = int(new_year)
        if new_sem: current_user.semester = int(new_sem)

        req_campus = request.form.get('campus')
        req_prog_id = request.form.get('program_id')

        if (req_campus and req_campus != current_user.campus) or \
           (req_prog_id and int(req_prog_id) != current_user.program_id):
            existing = TransferRequest.query.filter_by(user_id=current_user.id, status='pending').first()
            if not existing:
                new_req = TransferRequest(
                    user_id=current_user.id,
                    target_campus=req_campus or current_user.campus,
                    target_program_id=int(req_prog_id) if req_prog_id else current_user.program_id
                )
                db.session.add(new_req)
                flash('Institutional change request sent to Admin.', 'info')
                create_log("Transfer Request", f"User requested transfer to {req_campus} / Program {req_prog_id}")

        db.session.commit()
        create_log("Profile Update", f"Updated profile info (name changed from {old_name})")
        flash('Profile updated successfully!')
        return redirect(url_for('core.profile'))

    programs = Program.query.all()
    pending_request = TransferRequest.query.filter_by(user_id=current_user.id, status='pending').first()
    return render_template('core/student/profile.html', user=current_user, programs=programs,
                           pending_request=pending_request)





# ─── NOTES & RESOURCES ───
@core.route('/campus_notes')
@login_required
def campus_notes():
    notes = Resource.query.filter_by(campus=current_user.campus) \
        .join(Resource.programs) \
        .filter(Program.id == current_user.program_id) \
        .order_by(Resource.id.desc()) \
        .all()
    grouped_notes = {}
    for note in notes:
        m_name = note.module.name
        if m_name not in grouped_notes:
            grouped_notes[m_name] = []
        grouped_notes[m_name].append(note)
    return render_template('core/student/campus_notes.html', grouped_notes=grouped_notes)







@core.route('/cross_campus')
@login_required
def cross_campus():
    cross_notes = Resource.query.filter(Resource.campus != current_user.campus) \
        .join(Resource.programs) \
        .filter(Program.id == current_user.program_id) \
        .order_by(Resource.id.desc()) \
        .all()
    organized_notes = {}
    for note in cross_notes:
        campus = note.campus
        module_name = note.module.name
        if campus not in organized_notes:
            organized_notes[campus] = {}
        if module_name not in organized_notes[campus]:
            organized_notes[campus][module_name] = []
        organized_notes[campus][module_name].append(note)
    return render_template('core/student/cross_campus.html', organized_notes=organized_notes)






@core.route('/vault')
@login_required
def vault():
    return render_template('core/student/vault.html')




ALLOWED_EXTENSIONS = {'pdf', 'pptx', 'docx', 'doc'}

def generate_thumbnail(file_path, file_ext):
    """Generate thumbnail image data for the first page/slide of a document."""
    try:
        if file_ext == 'pdf':
            doc = fitz.open(file_path)
            page = doc[0]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            thumbnail_data = pix.tobytes("png")
            doc.close()
            return thumbnail_data

        elif file_ext in ['pptx', 'ppt']:
            prs = Presentation(file_path)
            if len(prs.slides) > 0:
                try:
                    title = prs.slides[0].shapes.title.text if prs.slides[0].shapes.title else "Slide 1"
                except:
                    title = "Slide 1"
                img = Image.new('RGB', (400, 300), color=(245, 242, 236))
                d = ImageDraw.Draw(img)
                try:
                    font = ImageFont.truetype("arial.ttf", 20)
                except:
                    font = ImageFont.load_default()
                d.text((20, 130), title, fill=(26, 26, 26), font=font)
                img_buffer = io.BytesIO()
                img.save(img_buffer, format='PNG')
                return img_buffer.getvalue()
            else:
                return None

        elif file_ext in ['docx', 'doc']:
            doc = Document(file_path)
            first_text = ""
            for para in doc.paragraphs:
                if para.text.strip():
                    first_text = para.text.strip()
                    break
            if not first_text:
                first_text = "Document Preview"
            img = Image.new('RGB', (400, 300), color=(245, 242, 236))
            d = ImageDraw.Draw(img)
            try:
                font = ImageFont.truetype("arial.ttf", 20)
            except:
                font = ImageFont.load_default()
            # Wrap text
            lines = []
            words = first_text.split()
            line = ""
            for word in words:
                if len(line + word) < 30:
                    line += word + " "
                else:
                    lines.append(line.strip())
                    line = word + " "
            if line:
                lines.append(line.strip())
            y = 130 - (len(lines) * 15)
            for l in lines:
                d.text((20, y), l, fill=(26, 26, 26), font=font)
                y += 30
            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            return img_buffer.getvalue()

    except Exception as e:
        print(f"Thumbnail generation failed: {e}")
        return None
    return None



@core.route('/upload-resource', methods=['GET', 'POST'])
@login_required
def upload_resource():
    if not current_user.is_class_rep:
        flash('Access denied.')
        return redirect(url_for('core.dashboard'))

    if request.method == 'POST':
        module_id = request.form.get('module_id')
        title = request.form.get('title')
        selected_programs = request.form.getlist('programs')
        file = request.files.get('file')

        if not module_id or not title or not selected_programs or not file:
            flash('All fields are required.', 'error')
            return redirect(url_for('core.upload_resource'))

        # Validate file type
        filename = secure_filename(file.filename)
        ext = filename.rsplit('.', 1)[-1].lower()
        if ext not in ALLOWED_EXTENSIONS:
            flash('Only PDF, PPTX, DOCX, and DOC files are allowed.', 'error')
            return redirect(url_for('core.upload_resource'))

        upload_folder = os.path.join(current_app.root_path, 'static', 'uploads')
        os.makedirs(upload_folder, exist_ok=True)
        file_path = os.path.join(upload_folder, filename)
        file.save(file_path)

        # Generate thumbnail
        thumbnail_data = generate_thumbnail(file_path, ext)
        thumbnail_filename = None
        if thumbnail_data:
            thumb_ext = 'png'
            thumbnail_filename = f"{os.path.splitext(filename)[0]}_thumb.{thumb_ext}"
            thumb_full_path = os.path.join(upload_folder, thumbnail_filename)
            with open(thumb_full_path, 'wb') as f:
                f.write(thumbnail_data)
            thumbnail_path = f"uploads/{thumbnail_filename}"
        else:
            thumbnail_path = None

        # Auto-fill academic_year and semester from rep's own data
        academic_year = current_user.year
        target_semester = current_user.semester

        new_resource = Resource(
            title=title,
            file_path=f"uploads/{filename}",
            thumbnail_path=thumbnail_path,
            file_type=ext,
            module_id=module_id,
            campus=current_user.campus,
            academic_year=academic_year,
            uploader_id=current_user.id,
            target_year=current_user.year,
            target_semester=current_user.semester
        )

        for prog_id in selected_programs:
            program = Program.query.get(int(prog_id))
            if program:
                new_resource.programs.append(program)

        db.session.add(new_resource)
        db.session.commit()

        # ─── SocketIO: Notify students in the same class ───
        room = f'student-{current_user.program_id}-{current_user.year}-{current_user.semester}'
        socketio.emit('new_resource', {
            'title': new_resource.title,
            'module_name': new_resource.module.name,
            'academic_year': new_resource.academic_year,
            'uploader_name': current_user.full_name,
            'campus': current_user.campus,
            'file_path': new_resource.file_path,
            'thumbnail_path': new_resource.thumbnail_path,
            'file_type': ext,
            'program_id': current_user.program_id,
            'year': current_user.year,
            'semester': current_user.semester
        }, room=room)

        # Also notify admins
        socketio.emit('new_resource', {
            'title': new_resource.title,
            'module_name': new_resource.module.name,
            'uploader_name': current_user.full_name,
            'campus': current_user.campus,
            'file_path': new_resource.file_path,
            'thumbnail_path': new_resource.thumbnail_path,
            'file_type': ext,
        }, room='admin')

        module = Module.query.get(module_id)
        create_log("Resource Upload", f"Uploaded '{title}' for module '{module.name}' (ID: {module.id})")
        flash('Resource uploaded successfully!')
        return redirect(url_for('core.dashboard'))

    modules = current_user.program.modules.all()
    programs = Program.query.all()
    return render_template('core/class_rep/upload_resource.html', modules=modules, programs=programs)





# ─── CLASS REP ACTIONS ───
@core.route('/approve-student/<int:user_id>', methods=['POST'])
@login_required
def approve_student(user_id):
    if not current_user.is_class_rep:
        flash('Access denied.')
        return redirect(url_for('core.dashboard'))

    student = User.query.get_or_404(user_id)
    if student.campus == current_user.campus and student.program_id == current_user.program_id:
        student.is_approved = True
        db.session.commit()

        # ─── SocketIO: Notify admin, remove from rep's pending list, inform student ───
        socketio.emit('student_approved', {
            'user_id': student.id,
            'full_name': student.full_name
        }, room='admin')

        rep_room = f'rep-{current_user.program_id}-{current_user.year}-{current_user.semester}'
        socketio.emit('pending_approval_removed', {
            'user_id': student.id
        }, room=rep_room)

        socketio.emit('approval_notification', {
            'message': 'Your account has been approved! You can now log in.'
        }, room=f'user-{student.id}')

        create_log("Student Approval", f"Approved student: {student.full_name} (ID: {student.id})")
        flash(f'Successfully approved {student.full_name}.')
    else:
        flash('You do not have permission to approve this student.')
    return redirect(url_for('core.dashboard'))





@core.route('/reject-student/<int:user_id>', methods=['POST'])
@login_required
def reject_student(user_id):
    if not current_user.is_class_rep:
        flash('Access denied.')
        return redirect(url_for('core.dashboard'))

    student = User.query.get_or_404(user_id)
    if student.campus == current_user.campus and student.program_id == current_user.program_id:
        student_name = student.full_name
        db.session.delete(student)
        db.session.commit()
        create_log("Registration Rejected", f"Rejected request from: {student_name} (ID: {user_id})")
        flash(f'Ignored registration request from {student_name}.')
    else:
        flash('You do not have permission to reject this student.')
    return redirect(url_for('core.dashboard'))





# ─── ADMIN DASHBOARD ───
def get_filtered_logs(action_filter=None, user_search=''):
    log_query = Log.query
    if action_filter:
        log_query = log_query.filter(Log.action == action_filter)

    if user_search:
        users_matching = User.query.filter(
            (User.full_name.ilike(f'%{user_search}%')) | (User.email.ilike(f'%{user_search}%'))
        ).all()
        user_ids = [u.id for u in users_matching]
        if user_ids:
            log_query = log_query.filter(Log.user_id.in_(user_ids))
        else:
            return []

    return log_query.order_by(Log.created_at.desc()).limit(50).all()





@core.route('/admin/logs')
@login_required
def admin_logs():
    if not current_user.is_admin:
        flash('Access denied.')
        return redirect(url_for('core.dashboard'))

    action_filter = request.args.get('action')
    user_search = request.args.get('user_search', '').strip()
    system_logs = get_filtered_logs(action_filter, user_search)

    return render_template('core/admin/_logs_table.html', logs=system_logs)






@core.route('/admin/dashboard')
@login_required
def admin_dashboard():
    if not current_user.is_admin:
        flash('Access denied.')
        return redirect(url_for('core.dashboard'))

    # Clean old logs (older than today)
    today = date.today()
    Log.query.filter(Log.created_at < today).delete()
    db.session.commit()

    # --- Filters ---
    year_filter = request.args.get('year', type=int)
    sem_filter = request.args.get('semester', type=int)
    action_filter = request.args.get('action')
    user_search = request.args.get('user_search', '').strip()

    # Pagination state for dashboard table previews
    page_reps = request.args.get('page_reps', 1, type=int)
    page_students = request.args.get('page_students', 1, type=int)
    page_modules = request.args.get('page_modules', 1, type=int)
    page_users = request.args.get('page_users', 1, type=int)
    page_logs = request.args.get('page_logs', 1, type=int)
    per_page = 10

    # 1. Reps
    reps_query = User.query.filter_by(is_class_rep=True, is_admin=False)
    reps_paginated = paginate_query(reps_query, page_reps, per_page)

    # 2. Students
    student_query = User.query.filter(User.is_class_rep == False, User.is_admin == False)
    if year_filter:
        student_query = student_query.filter_by(year=year_filter)
    if sem_filter:
        student_query = student_query.filter_by(semester=sem_filter)
    students_paginated = paginate_query(student_query, page_students, per_page)

    # 3. Modules
    modules_query = Module.query
    modules_paginated = paginate_query(modules_query, page_modules, per_page)

    # 4. Users
    users_query = User.query
    users_paginated = paginate_query(users_query, page_users, per_page)

    # 5. Logs
    log_query = Log.query
    if action_filter:
        log_query = log_query.filter(Log.action == action_filter)
    if user_search:
        users_matching = User.query.filter(
            (User.full_name.ilike(f'%{user_search}%')) | (User.email.ilike(f'%{user_search}%'))
        ).all()
        user_ids = [u.id for u in users_matching]
        if user_ids:
            log_query = log_query.filter(Log.user_id.in_(user_ids))
        else:
            log_query = log_query.filter(False)
    logs_paginated = paginate_query(log_query.order_by(Log.created_at.desc()), page_logs, per_page)

    # 3. Get distinct actions for dropdown
    distinct_actions = db.session.query(Log.action).distinct().all()
    actions_list = [a[0] for a in distinct_actions if a[0] is not None]

    # 4. Current time for online status
    now = datetime.utcnow()

    # 5. Other data
    all_programs = Program.query.all()
    all_modules = Module.query.all()
    pending_requests = TransferRequest.query.filter_by(status='pending').all()

    # Log admin dashboard access (only if filter applied, to avoid clutter)
    if year_filter or sem_filter or action_filter or user_search:
        create_log("Admin Dashboard Filter", f"Filtered by Year={year_filter}, Semester={sem_filter}, Action={action_filter}, User={user_search}")

    return render_template('core/admin/dashboard.html',
                           reps_users=reps_paginated['items'],
                           page_reps=reps_paginated['page'],
                           pages_reps=reps_paginated['pages'],
                           student_users=students_paginated['items'],
                           page_students=students_paginated['page'],
                           pages_students=students_paginated['pages'],
                           modules=modules_paginated['items'],
                           page_modules=modules_paginated['page'],
                           pages_modules=modules_paginated['pages'],
                           admin_users=users_paginated['items'],
                           page_users=users_paginated['page'],
                           pages_users=users_paginated['pages'],
                           log_items=logs_paginated['items'],
                           page_logs=logs_paginated['page'],
                           pages_logs=logs_paginated['pages'],
                           programs=all_programs,
                           modules_all=all_modules,
                           requests=pending_requests,
                           current_year=year_filter,
                           current_sem=sem_filter,
                           actions_list=actions_list,
                           selected_action=action_filter,
                           user_search=user_search,
                           now=now)





@core.route('/live/admin-logs')
@login_required
def live_admin_logs():
    if not current_user.is_admin:
        return jsonify({'error': 'Unauthorized'}), 401

    action_filter = request.args.get('action')
    user_search = request.args.get('user_search', '').strip()
    system_logs = get_filtered_logs(action_filter, user_search)

    logs_data = []
    for log in system_logs:
        logs_data.append({
            'created_at': log.created_at.strftime('%H:%M | %d %b %Y') if log.created_at else '',
            'actor': log.user.full_name if log.user else 'System Core',
            'action': log.action or '',
            'details': log.details or ''
        })

    return jsonify({
        'logs': logs_data,
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    })





@core.route('/admin/request/<int:req_id>/<action>')
@login_required
def handle_request(req_id, action):
    if not current_user.is_admin:
        abort(403)

    req = TransferRequest.query.get_or_404(req_id)
    if action == 'approve':
        req.user.campus = req.target_campus
        req.user.program_id = req.target_program_id
        req.status = 'approved'
        create_log("Transfer Approved", f"Approved transfer for {req.user.full_name} (ID: {req.user.id})")
        flash(f"Approved changes for {req.user.full_name}")
    else:
        req.status = 'rejected'
        create_log("Transfer Rejected", f"Rejected transfer for {req.user.full_name} (ID: {req.user.id})")
        flash("Request rejected.")

    db.session.commit()
    return redirect(url_for('core.admin_dashboard'))





@core.route('/admin/add-module', methods=['POST'])
@login_required
def add_module():
    if not current_user.is_admin:
        flash('Access denied.')
        return redirect(url_for('core.dashboard'))

    name = request.form.get('name', '').strip()
    code = request.form.get('code', '').strip().upper()
    program_ids = request.form.getlist('programs')

    if not name or not code or not program_ids:
        flash('Error: Module Name, Code, and at least one Program are required!')
        return redirect(url_for('core.admin_dashboard'))

    existing = Module.query.filter_by(code=code).first()
    if existing:
        flash(f'Error: Module code "{code}" is already used by "{existing.name}".')
        return redirect(url_for('core.admin_dashboard'))

    try:
        new_module = Module(name=name, code=code)
        for p_id in program_ids:
            prog = Program.query.get(int(p_id))
            if prog:
                new_module.programs.append(prog)
        db.session.add(new_module)
        db.session.commit()

        # ─── SocketIO: Notify all admins ───
        socketio.emit('module_added', {
            'module_id': new_module.id,
            'name': new_module.name,
            'code': new_module.code,
            'programs': [p.name for p in new_module.programs]
        }, room='admin')

        create_log("Module Created", f"Admin registered module: {name} ({code}) linked to {len(program_ids)} programs")
        flash(f'Module {name} registered successfully!')
    except Exception as e:
        db.session.rollback()
        flash('Something went wrong.')
        print(f"Database Error: {e}")
    return redirect(url_for('core.admin_dashboard'))




@core.route('/admin/make-rep/<int:user_id>', methods=['POST'])
@login_required
def make_rep(user_id):
    if not current_user.is_admin:
        flash('Access denied.', 'error')
        return redirect(url_for('core.dashboard'))

    user = User.query.get_or_404(user_id)
    user.is_class_rep = True
    user.is_approved = True
    db.session.commit()

    # ─── SocketIO: Notify admin and user ───
    socketio.emit('user_role_changed', {
        'user_id': user.id,
        'full_name': user.full_name,
        'new_role': 'rep',
        'is_approved': user.is_approved
    }, room='admin')

    socketio.emit('role_update', {
        'new_role': 'rep'
    }, room=f'user-{user.id}')

    create_log("Role Change", f"Promoted and approved {user.full_name} (ID: {user.id}) as Class Representative")
    flash(f'{user.full_name} is now a Class Representative and can log in!', 'success')
    return redirect(url_for('core.admin_dashboard'))





@core.route('/admin/remove-rep/<int:user_id>', methods=['POST'])
@login_required
def remove_rep(user_id):
    if not current_user.is_admin:
        flash('Access denied.', 'error')
        return redirect(url_for('core.dashboard'))

    user = User.query.get_or_404(user_id)
    user.is_class_rep = False
    db.session.commit()
    create_log("Role Change", f"Removed {user.full_name} (ID: {user.id}) from being Class Representative")
    flash(f'{user.full_name} is no longer a Class Representative!', 'info')
    return redirect(url_for('core.admin_dashboard'))




@core.route('/admin/delete-module/<int:module_id>', methods=['POST'])
@login_required
def delete_module(module_id):
    if not current_user.is_admin:
        flash('Access denied.', 'error')
        return redirect(url_for('core.dashboard'))

    module = Module.query.get_or_404(module_id)
    if module.resources:
        flash(f'Cannot delete "{module.name}" because it contains resources. Delete the notes first!', 'error')
        return redirect(url_for('core.admin_dashboard'))

    module_name = module.name
    db.session.delete(module)
    db.session.commit()
    create_log("Module Deleted", f"Deleted Module: {module_name} (ID: {module_id})")
    flash(f'Module "{module_name}" removed successfully.', 'info')
    return redirect(url_for('core.admin_dashboard'))




@core.route('/notes/delete/<int:note_id>', methods=['POST'])
@login_required
def delete_note(note_id):
    note = Resource.query.get_or_404(note_id)
    if not current_user.is_class_rep or current_user.id != note.user_id:
        flash('Permission denied. You can only delete your own uploads.', 'error')
        return redirect(url_for('core.campus_notes'))

    try:
        file_path = os.path.join(current_app.root_path, 'static', note.file_path)
        if os.path.exists(file_path):
            os.remove(file_path)
        db.session.delete(note)
        db.session.commit()
        create_log("File Removed", f"Rep {current_user.full_name} deleted note: {note.title} (ID: {note_id})")
        flash('File deleted successfully.', 'info')
    except Exception as e:
        db.session.rollback()
        print(f"Delete Error: {e}")
        flash('Something went wrong while removing the file.', 'error')
    return redirect(url_for('core.campus_notes'))




# ─── ADMIN CHANGE PASSWORD ───
@core.route('/admin/change-password/<int:user_id>', methods=['POST'])
@login_required
def admin_change_password(user_id):
    if not current_user.is_admin:
        flash('Access denied.', 'error')
        return redirect(url_for('core.admin_dashboard'))

    user = User.query.get_or_404(user_id)
    new_password = request.form.get('new_password')
    confirm_password = request.form.get('confirm_password')

    if not new_password or len(new_password) < 6:
        flash('Password must be at least 6 characters.', 'error')
        return redirect(url_for('core.admin_dashboard'))

    if new_password != confirm_password:
        flash('Passwords do not match.', 'error')
        return redirect(url_for('core.admin_dashboard'))

    user.password_hash = generate_password_hash(new_password)
    db.session.commit()
    create_log("Password Reset", f"Admin reset password for {user.full_name} (ID: {user.id})")
    flash(f'Password for {user.full_name} has been reset successfully.', 'success')
    return redirect(url_for('core.admin_dashboard'))





# ─── SETTINGS ───
@core.route('/settings', methods=['GET', 'POST'])
@login_required
def settings():
    if request.method == 'POST':
        theme = request.form.get('theme')
        lang = request.form.get('language')

        if theme in ['light', 'dark']:
            current_user.theme_preference = theme
        if lang in ['en']:
            current_user.language_preference = lang

        db.session.commit()
        create_log("Settings Update", f"User updated theme to '{theme}', language to '{lang}'")
        flash('Settings updated successfully!', 'success')
        return redirect(url_for('core.settings'))

    return render_template('core/student/settings.html', user=current_user)




@core.route('/settings/deactivate', methods=['POST'])
@login_required
def deactivate_account():
    if current_user.is_admin:
        flash('Admin accounts cannot be deactivated via this form. Contact another admin.', 'error')
        return redirect(url_for('core.settings'))

    current_user.is_active = False
    db.session.commit()
    create_log("Account Deactivation", f"User {current_user.full_name} (ID: {current_user.id}) deactivated their account")
    logout_user()
    flash('Your account has been deactivated. You can contact an admin to reactivate.', 'info')
    return redirect(url_for('auth.login'))






# ─── ANALYTICS ───
@core.route('/admin/analytics')
@login_required
def admin_analytics():
    if not current_user.is_admin:
        flash('Access denied.', 'error')
        return redirect(url_for('core.dashboard'))

    # Students per module
    module_counts = db.session.query(
        Module.name,
        func.count(User.id).label('count')
    ).join(Module.programs).join(Program.users).filter(User.is_approved == True).group_by(Module.id).all()
    modules_labels = [row[0] for row in module_counts] or ['No modules']
    modules_data = [row[1] for row in module_counts] or [0]

    # Students per year
    year_counts = db.session.query(
        User.year,
        func.count(User.id).label('count')
    ).filter(User.is_approved == True).group_by(User.year).all()
    years_dict = {y: 0 for y in range(1, 5)}
    for y, c in year_counts:
        years_dict[y] = c
    years_labels = ['Year 1', 'Year 2', 'Year 3', 'Year 4']
    years_data = [years_dict[1], years_dict[2], years_dict[3], years_dict[4]]

    # Students per semester
    sem_counts = db.session.query(
        User.semester,
        func.count(User.id).label('count')
    ).filter(User.is_approved == True).group_by(User.semester).all()
    sem_dict = {1: 0, 2: 0}
    for s, c in sem_counts:
        sem_dict[s] = c
    sem_labels = ['Semester 1', 'Semester 2']
    sem_data = [sem_dict[1], sem_dict[2]]

    # Students per campus
    campus_counts = db.session.query(
        User.campus,
        func.count(User.id).label('count')
    ).filter(User.is_approved == True).group_by(User.campus).all()
    campus_labels = [row[0] for row in campus_counts]
    campus_data = [row[1] for row in campus_counts]

    # Approval status
    approved_count = User.query.filter_by(is_approved=True).count()
    pending_count = User.query.filter_by(is_approved=False).count()
    status_labels = ['Approved', 'Pending']
    status_data = [approved_count, pending_count]

    return render_template('core/admin/analytics.html',
                           modules_labels=modules_labels,
                           modules_data=modules_data,
                           years_labels=years_labels,
                           years_data=years_data,
                           sem_labels=sem_labels,
                           sem_data=sem_data,
                           campus_labels=campus_labels,
                           campus_data=campus_data,
                           status_labels=status_labels,
                           status_data=status_data)





# ─── API ENDPOINTS ───
@core.route('/api/login', methods=['POST'])
def api_login():
    data = request.json
    email = data.get('email')
    password = data.get('password')

    user = User.query.filter_by(email=email).first()

    if user and check_password_hash(user.password_hash, password):
        if user.is_approved:
            login_user(user, remember=True)
            response = jsonify({"status": "success", "message": "Welcome!"})
            response.headers.add("ngrok-skip-browser-warning", "true")
            response.set_cookie('session', request.cookies.get('session', ''), samesite='None', secure=True)
            return response, 200
        else:
            return jsonify({"status": "error", "message": "Account pending approval"}), 403
    else:
        return jsonify({"status": "error", "message": "Invalid credentials"}), 401





@core.route('/api/v1/student-context', methods=['GET'])
@login_required
def get_student_context():
    active_modules = current_user.program.modules.all()
    modules_list = [{"name": m.name, "code": m.code} for m in active_modules]
    response = jsonify({
        "student_name": current_user.full_name,
        "modules": modules_list
    })
    response.headers.add("ngrok-skip-browser-warning", "true")
    return response

@core.route('/about')
def about():
    return render_template('core/about.html')